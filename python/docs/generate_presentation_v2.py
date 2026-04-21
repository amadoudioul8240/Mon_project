from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parents[2]
DOCS_DIR = ROOT / "docs"
DATE_STR = datetime.now().strftime("%Y-%m-%d")
OUTPUT = DOCS_DIR / f"IT-Monitoring-Presentation-V2-{DATE_STR}.pptx"


PRIMARY = RGBColor(11, 61, 145)
ACCENT = RGBColor(10, 147, 150)
DARK = RGBColor(20, 24, 28)
LIGHT = RGBColor(245, 248, 252)
WHITE = RGBColor(255, 255, 255)


def add_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()

    top_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.7))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = PRIMARY
    top_band.line.fill.background()


def add_title(slide, title, subtitle=""):
    add_background(slide)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.9), Inches(12.0), Inches(1.2))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = DARK

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(11.5), Inches(0.8))
        stf = sub_box.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(18)
        sp.font.color.rgb = RGBColor(70, 80, 90)


def add_kpi_card(slide, left, top, w, h, value, label):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(210, 220, 230)

    vbox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(w - 0.4), Inches(0.6))
    vtf = vbox.text_frame
    vp = vtf.paragraphs[0]
    vp.text = value
    vp.font.bold = True
    vp.font.size = Pt(28)
    vp.font.color.rgb = PRIMARY

    lbox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.95), Inches(w - 0.4), Inches(0.55))
    ltf = lbox.text_frame
    lp = ltf.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(14)
    lp.font.color.rgb = RGBColor(60, 70, 80)


def add_bullet_block(slide, x, y, w, h, title, bullets):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(210, 220, 230)

    tbox = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.15), Inches(w - 0.5), Inches(0.5))
    ttf = tbox.text_frame
    tp = ttf.paragraphs[0]
    tp.text = title
    tp.font.bold = True
    tp.font.size = Pt(18)
    tp.font.color.rgb = PRIMARY

    bbox = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.65), Inches(w - 0.6), Inches(h - 0.8))
    btf = bbox.text_frame
    btf.clear()
    for i, b in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(15)
        p.font.color.rgb = DARK


def add_footer(slide, text="IT Monitoring - Presentation V2"):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.1), Inches(8.0), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(90, 100, 110)


def build_presentation(path: Path):
    prs = Presentation()

    # 1. Cover
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s1, "IT Monitoring", "Soutenance professionnelle - supervision et gestion du parc IT")
    add_kpi_card(s1, 0.8, 3.2, 3.0, 1.8, "3", "Composants coeur")
    add_kpi_card(s1, 4.15, 3.2, 3.0, 1.8, "59", "Endpoints API")
    add_kpi_card(s1, 7.5, 3.2, 3.0, 1.8, "13+", "Pages metier")
    add_footer(s1)

    # 2. Problem and objectives
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s2, "Pourquoi ce projet ?", "Unifier la visibilite IT et accelerer l'operationnel")
    add_bullet_block(s2, 0.8, 2.2, 5.9, 3.8, "Problemes addresses", [
        "Inventaire eparpille et peu fiable",
        "Faible visibilite securite endpoint",
        "Difficulte de priorisation incidents",
        "Collecte terrain non industrialisee",
    ])
    add_bullet_block(s2, 6.9, 2.2, 5.6, 3.8, "Objectifs", [
        "Centraliser les donnees techniques",
        "Suivre actifs, incidents, projets",
        "Mesurer posture et risques en continu",
        "Poser une base SIEM operationnelle",
    ])
    add_footer(s2)

    # 3. Architecture
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s3, "Architecture de reference")

    blocks = [
        (0.9, 2.1, 2.8, 1.3, "Frontend", "React"),
        (4.0, 2.1, 2.8, 1.3, "Backend", "FastAPI"),
        (7.1, 2.1, 2.8, 1.3, "Database", "PostgreSQL"),
        (10.2, 2.1, 2.2, 1.3, "Agents", "Go/PS/Py"),
    ]
    for x, y, w, h, name, tech in blocks:
        shape = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = ACCENT
        tx = s3.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.15), Inches(w - 0.3), Inches(0.9))
        tf = tx.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = name
        p1.font.bold = True
        p1.font.size = Pt(18)
        p1.font.color.rgb = PRIMARY
        p2 = tf.add_paragraph()
        p2.text = tech
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK

    add_bullet_block(s3, 0.9, 4.1, 11.5, 2.1, "Flux principal", [
        "Agents -> Backend (ingestion) -> PostgreSQL -> Frontend dashboards",
        "Docker Compose pour orchestration locale et serveur",
        "CORS et URL backend adaptes au serveur 192.168.196.134",
    ])
    add_footer(s3)

    # 4. Functional modules
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s4, "Modules fonctionnels")
    add_bullet_block(s4, 0.8, 2.0, 4.0, 4.8, "Pilotage parc", [
        "Assets et types",
        "Utilisateurs et sites",
        "Maintenance logs",
        "Statistiques inventaire",
    ])
    add_bullet_block(s4, 4.95, 2.0, 4.0, 4.8, "Operations", [
        "Incidents",
        "Projets IT",
        "Supervision ressources",
        "Ports et logs LAN",
    ])
    add_bullet_block(s4, 9.1, 2.0, 3.4, 4.8, "Securite", [
        "Findings et policy",
        "Posture endpoint",
        "CVE watch",
        "SIEM + timeline",
    ])
    add_footer(s4)

    # 5. API
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s5, "API platforme")
    add_bullet_block(s5, 0.9, 2.0, 5.8, 2.2, "Domaines", [
        "Actifs, referentiels, maintenance",
        "AD sync, securite, SIEM, reseau, metriques",
    ])
    add_bullet_block(s5, 7.0, 2.0, 5.4, 2.2, "Volume", [
        "59 endpoints FastAPI", "/docs disponible pour exploration",
    ])
    add_bullet_block(s5, 0.9, 4.5, 11.5, 2.2, "Exemples critiques", [
        "/assets/scan, /metrics/resources, /security/summary, /siem/alerts, /siem/timeline/{host_serial}",
    ])
    add_footer(s5)

    # 6. Agents
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s6, "Agents de collecte terrain")
    add_bullet_block(s6, 0.8, 2.0, 5.8, 4.8, "Capacites", [
        "Inventaire machine et logiciels",
        "Posture securite endpoint",
        "Metriques CPU/RAM/disque",
        "Ports ouverts et telemetrie LAN",
        "Evenements d'authentification",
    ])
    add_bullet_block(s6, 6.9, 2.0, 5.6, 4.8, "Deploiement", [
        "Service Windows (Go)",
        "Modes: manuel, bulk, GPO",
        "ASSET_BACKEND_URL surchargeable",
        "Defaut: 192.168.196.134:8000/assets/scan",
    ])
    add_footer(s6)

    # 7. Deployment and ops
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s7, "Deploiement et exploitation")
    add_bullet_block(s7, 0.9, 2.0, 11.5, 4.8, "Runbook serveur Debian", [
        "git pull origin main dans le repo cible",
        "docker compose up -d --build",
        "Verification ports 3000/8000 et endpoints HTTP",
        "Diagnostic par logs conteneurs backend/frontend",
        "Gestion CORS via CORS_ALLOW_ORIGINS",
    ])
    add_footer(s7)

    # 8. Security and risks
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s8, "Securite, risques, maitrise")
    add_bullet_block(s8, 0.8, 2.0, 5.8, 4.8, "Maitrise actuelle", [
        "Collecte securite centralisee",
        "SIEM interne avec alerting",
        "Signature agent possible (all-signed)",
        "Configuration CORS controlee",
    ])
    add_bullet_block(s8, 6.9, 2.0, 5.6, 4.8, "Points d'attention", [
        "Industrialiser tests CI/CD",
        "Renforcer backup/restauration DB",
        "Generaliser TLS/reverse proxy",
        "Suivre KPI MTTD/MTTR",
    ])
    add_footer(s8)

    # 9. Roadmap
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s9, "Roadmap 90 jours")
    add_kpi_card(s9, 0.9, 2.2, 3.8, 1.9, "J+30", "Stabilisation prod")
    add_kpi_card(s9, 4.9, 2.2, 3.8, 1.9, "J+60", "Automatisation CI/CD")
    add_kpi_card(s9, 8.9, 2.2, 3.8, 1.9, "J+90", "KPI exec et compliance")
    add_bullet_block(s9, 0.9, 4.5, 11.8, 2.1, "Priorites", [
        "Fiabilite deploiement, observabilite, gouvernance securite, scalabilite fonctionnelle",
    ])
    add_footer(s9)

    # 10. Closing
    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s10, "Conclusion", "Une base solide pour une plateforme IT operationnelle et evolutive")
    add_bullet_block(s10, 1.3, 2.4, 10.8, 3.6, "Messages clefs", [
        "Architecture moderne, modulaire et exploitable",
        "Couverture metier large: inventaire, ops, securite, SIEM",
        "Agents adaptes au terrain et redeployables facilement",
        "Projet pret pour une presentation professionnelle et une phase d'industrialisation",
    ])
    add_footer(s10, "Merci - IT Monitoring")

    prs.save(str(path))


if __name__ == "__main__":
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    build_presentation(OUTPUT)
    print(f"PPTX generated: {OUTPUT}")
