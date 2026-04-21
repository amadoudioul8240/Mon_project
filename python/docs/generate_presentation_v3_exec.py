from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parents[2]
DOCS_DIR = ROOT / "docs"
DATE_STR = datetime.now().strftime("%Y-%m-%d")
OUTPUT = DOCS_DIR / f"IT-Monitoring-Presentation-V3-Executive-{DATE_STR}.pptx"


PRIMARY = RGBColor(12, 43, 104)
ACCENT = RGBColor(0, 122, 204)
SUCCESS = RGBColor(28, 126, 62)
WARNING = RGBColor(188, 124, 16)
DARK = RGBColor(24, 28, 34)
LIGHT = RGBColor(246, 249, 253)
WHITE = RGBColor(255, 255, 255)


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.58))
    top.fill.solid()
    top.fill.fore_color.rgb = PRIMARY
    top.line.fill.background()


def add_title(slide, title, subtitle=None):
    add_bg(slide)

    tbox = slide.shapes.add_textbox(Inches(0.65), Inches(0.8), Inches(12.2), Inches(0.9))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(34)
    p.font.color.rgb = DARK

    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.67), Inches(1.62), Inches(12.0), Inches(0.6))
        stf = sbox.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(17)
        sp.font.color.rgb = RGBColor(85, 96, 110)


def add_footer(slide):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.12), Inches(12), Inches(0.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"IT Monitoring | Executive deck | {DATE_STR}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(110, 120, 130)


def add_card(slide, x, y, w, h, title, value, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(215, 223, 232)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.14))
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()

    t = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.25), Inches(w - 0.4), Inches(0.45))
    tp = t.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(13)
    tp.font.color.rgb = RGBColor(82, 92, 102)

    v = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.72), Inches(w - 0.4), Inches(0.6))
    vp = v.text_frame.paragraphs[0]
    vp.text = value
    vp.font.bold = True
    vp.font.size = Pt(25)
    vp.font.color.rgb = DARK


def add_bullets(slide, x, y, w, h, title, bullets):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(215, 223, 232)

    th = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.16), Inches(w - 0.4), Inches(0.45))
    tp = th.text_frame.paragraphs[0]
    tp.text = title
    tp.font.bold = True
    tp.font.size = Pt(17)
    tp.font.color.rgb = PRIMARY

    box = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.62), Inches(w - 0.45), Inches(h - 0.75))
    tf = box.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(14)
        p.font.color.rgb = DARK


def build(path: Path):
    prs = Presentation()

    # Slide 1 - Executive summary
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s1, "IT Monitoring", "Synthese executive pour comite de direction")
    add_card(s1, 0.75, 2.3, 3.0, 1.7, "Perimetre", "Parc IT unifie")
    add_card(s1, 3.95, 2.3, 3.0, 1.7, "Stack", "React + FastAPI")
    add_card(s1, 7.15, 2.3, 2.8, 1.7, "API", "59 endpoints")
    add_card(s1, 10.15, 2.3, 2.4, 1.7, "Cible", "Ops + Sec", SUCCESS)
    add_bullets(s1, 0.75, 4.35, 11.8, 2.3, "Message cle", [
        "Plateforme operationnelle qui centralise l'inventaire, la supervision et la securite endpoint.",
        "Base prete pour industrialisation CI/CD, gouvernance et reporting de performance.",
    ])
    add_footer(s1)

    # Slide 2 - Business value
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s2, "Valeur metier", "Pourquoi investir dans la phase d'industrialisation")
    add_bullets(s2, 0.75, 2.05, 5.75, 4.8, "Gains attendus", [
        "Reduction du temps de diagnostic incident",
        "Meilleure fiabilite de l'inventaire IT",
        "Visibilite securite centralisee et continue",
        "Decision plus rapide grace aux dashboards", 
    ])
    add_bullets(s2, 6.75, 2.05, 5.75, 4.8, "Impact organisationnel", [
        "Standardisation des processus N1/N2",
        "Montree en competence progressive equipe ops",
        "Traçabilite des actions securite",
        "Support a la conformite et audits",
    ])
    add_footer(s2)

    # Slide 3 - KPI dashboard
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s3, "KPI de pilotage", "Proposition de tableau de bord direction")
    add_card(s3, 0.75, 2.2, 2.85, 1.55, "Disponibilite plate-forme", ">99%", SUCCESS)
    add_card(s3, 3.8, 2.2, 2.85, 1.55, "Couverture endpoints", "Cible 85%", ACCENT)
    add_card(s3, 6.85, 2.2, 2.85, 1.55, "MTTD securite", "-30%", SUCCESS)
    add_card(s3, 9.9, 2.2, 2.65, 1.55, "MTTR incident", "-25%", SUCCESS)

    add_card(s3, 0.75, 3.95, 2.85, 1.55, "Taux actifs qualifies", ">95%", ACCENT)
    add_card(s3, 3.8, 3.95, 2.85, 1.55, "Alerts traitees <24h", ">90%", ACCENT)
    add_card(s3, 6.85, 3.95, 2.85, 1.55, "Derive capacitaire", "<10%", WARNING)
    add_card(s3, 9.9, 3.95, 2.65, 1.55, "Conformite endpoint", ">92%", ACCENT)

    add_bullets(s3, 0.75, 5.75, 11.8, 1.0, "Note", ["KPI a affiner avec les donnees de production sur les 4 premieres semaines."])
    add_footer(s3)

    # Slide 4 - Budget and ROI
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s4, "Budget et ROI", "Hypothese de cadrage pour phase 2 (90 jours)")
    add_bullets(s4, 0.75, 2.1, 5.75, 4.9, "Budget indicatif", [
        "Ingenierie/devops: 30-45 j.h",
        "Durcissement securite: 10-15 j.h",
        "Accompagnement exploitation: 8-12 j.h",
        "Formation et transfert: 5-8 j.h",
        "Total estime: 53-80 j.h",
    ])
    add_bullets(s4, 6.75, 2.1, 5.75, 4.9, "ROI attendu (6-12 mois)", [
        "Moins d'interruptions non maitrisees",
        "Baisse des couts de diagnostic manuel",
        "Gain de productivite equipes IT",
        "Reduction exposition risques securite",
        "Retour sur investissement probable <12 mois",
    ])
    add_footer(s4)

    # Slide 5 - Risks and mitigation
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s5, "Risques et plans de maitrise")
    add_bullets(s5, 0.75, 2.05, 5.75, 4.9, "Risques majeurs", [
        "Couverture agent incomplete",
        "Dette technique CI/CD et tests",
        "Dependance a des configurations manuelles",
        "Sensibilite DNS/reseau en environnement serveur",
    ])
    add_bullets(s5, 6.75, 2.05, 5.75, 4.9, "Actions de mitigation", [
        "Automatiser deploiement et controles post-deploy",
        "Mettre en place pipeline de tests obligatoires",
        "Durcir observabilite + alerting technique",
        "Documenter runbooks N1/N2 + PRA",
    ])
    add_footer(s5)

    # Slide 6 - Decision and next steps
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s6, "Decision attendue", "Validation du passage en phase industrialisation")
    add_bullets(s6, 0.9, 2.2, 11.6, 3.0, "Proposition", [
        "Lancer un sprint de 90 jours axe fiabilite, securite et pilotage KPI.",
        "Nommer un sponsor metier + un referent exploitation.",
        "Mettre en place un reporting hebdomadaire au comite de suivi.",
    ])

    add_card(s6, 1.2, 5.45, 3.6, 1.4, "J+30", "Stabilisation prod", SUCCESS)
    add_card(s6, 4.95, 5.45, 3.6, 1.4, "J+60", "Automatisation CI/CD", ACCENT)
    add_card(s6, 8.7, 5.45, 3.6, 1.4, "J+90", "Tableau de bord exec", PRIMARY)
    add_footer(s6)

    prs.save(str(path))


if __name__ == "__main__":
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    build(OUTPUT)
    print(f"PPTX generated: {OUTPUT}")
