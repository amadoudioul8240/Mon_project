from pathlib import Path
from datetime import datetime

from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle, PageBreak

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parents[2]
DOCS_DIR = ROOT / "docs"
DATE_STR = datetime.now().strftime("%Y-%m-%d")

PDF_PATH = DOCS_DIR / f"IT-Monitoring-Documentation-{DATE_STR}.pdf"
PPTX_PATH = DOCS_DIR / f"IT-Monitoring-Presentation-{DATE_STR}.pptx"


PROJECT_SUMMARY = [
    "IT Monitoring est une plateforme de gestion et de supervision du parc informatique.",
    "Le projet combine un backend FastAPI, un frontend React et une base PostgreSQL.",
    "La solution inclut des agents de collecte (PowerShell, Python, Go) pour la remontee terrain.",
    "Le deploiement est conteneurise avec Docker Compose pour simplifier l'exploitation.",
]

ARCH_COMPONENTS = [
    ["Composant", "Technologie", "Role"],
    ["Frontend", "React 18 + Recharts", "Interface operateur, dashboards, pages metier"],
    ["Backend", "FastAPI + SQLAlchemy", "API REST, logique metier, AD sync, SIEM interne"],
    ["Base de donnees", "PostgreSQL 15", "Persistances assets, incidents, telemetrie, SIEM"],
    ["Agent Windows", "Go service + scripts PS", "Collecte inventaire, securite, ressources, reseau"],
    ["Orchestration", "Docker Compose", "Execution et cycle de vie des services"],
]

KEY_FEATURES = [
    "Inventaire IT: utilisateurs, sites, actifs, types d'actifs.",
    "Cycle de vie actif: creation, mise a jour, suppression, maintenance logs.",
    "Synchronisation AD: statut, configuration, sync users/computers.",
    "Supervision ressources endpoint: CPU/RAM/disque + vue globale.",
    "Volet securite: findings, policy, posture endpoint, summary.",
    "Network telemetry: ports ouverts, logs LAN, appareils inconnus.",
    "Mini SIEM: collecte auth/download/web events, alerting et timeline.",
    "Pilotage projets IT et incidents pour operation quotidienne.",
]

API_DOMAINS = [
    ["Domaine API", "Exemples d'endpoints"],
    ["Actifs et referentiels", "/assets, /asset_types, /users, /locations"],
    ["Maintenance", "/maintenance_logs/{asset_id}"],
    ["Agent ingestion", "/assets/scan, /metrics/resources, /network/telemetry"],
    ["AD integration", "/ad/status, /ad/config, /ad/sync/users, /ad/sync/computers"],
    ["Securite", "/security/findings, /security/policy, /security/posture, /security/summary"],
    ["SIEM", "/siem/auth-events, /siem/alerts, /siem/timeline/{host_serial}"],
    ["Operations", "/incidents, /it-projects, /stats/assets-by-type"],
]

DEPLOY_STEPS = [
    "1) Cloner le projet sur le serveur Debian de deploiement.",
    "2) Configurer .env (si necessaire) et verifier docker compose.",
    "3) Lancer: docker compose up -d --build.",
    "4) Verifier: docker compose ps, curl http://IP:8000/docs, curl http://IP:3000.",
    "5) Supervision: logs conteneurs backend/frontend en cas d'anomalie.",
]

AGENT_NOTES = [
    "L'URL backend par defaut des agents pointe vers http://192.168.196.134:8000/assets/scan.",
    "La variable ASSET_BACKEND_URL permet de surcharger la cible sans modifier le code.",
    "Les modes de deploiement agent incluent: manuel (install_service), bulk, GPO, all-signed.",
]

SECURITY_PRACTICES = [
    "Signer les binaires/scripts agents pour execution en environnement durci.",
    "Limiter CORS aux origines frontend autorisees via CORS_ALLOW_ORIGINS.",
    "Utiliser comptes de service dedies pour AD sync avec moindre privilege.",
    "Journaliser les evenements critiques et monitorer les erreurs d'ingestion.",
]

ROADMAP = [
    "CI/CD mature avec tests automatiques backend/frontend avant deploiement.",
    "Alerting avance (email/Teams/Slack) sur incidents SIEM et derive ressources.",
    "Durcissement infra: TLS, reverse proxy, segmentation reseau, sauvegarde DB.",
    "Tableau de bord executif (KPI SLA, MTTD, MTTR, conformite endpoint).",
]



def header_footer(canvas, doc):
    canvas.saveState()
    canvas.setFont("Helvetica", 9)
    canvas.setFillColor(colors.grey)
    canvas.drawString(2 * cm, 1.2 * cm, f"IT Monitoring - Documentation professionnelle")
    canvas.drawRightString(A4[0] - 2 * cm, 1.2 * cm, f"Page {doc.page}")
    canvas.restoreState()



def build_pdf(path: Path) -> None:
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "DocTitle",
        parent=styles["Title"],
        fontName="Helvetica-Bold",
        fontSize=24,
        leading=30,
        textColor=colors.HexColor("#0B3D91"),
        alignment=1,
        spaceAfter=18,
    )
    h1 = ParagraphStyle(
        "H1",
        parent=styles["Heading1"],
        fontName="Helvetica-Bold",
        fontSize=16,
        leading=20,
        textColor=colors.HexColor("#12355B"),
        spaceBefore=10,
        spaceAfter=8,
    )
    body = ParagraphStyle(
        "Body",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=11,
        leading=16,
    )

    doc = SimpleDocTemplate(
        str(path),
        pagesize=A4,
        leftMargin=2 * cm,
        rightMargin=2 * cm,
        topMargin=2 * cm,
        bottomMargin=2.2 * cm,
        title="IT Monitoring - Documentation complete",
        author="IT Monitoring Team",
    )

    story = []
    story.append(Paragraph("IT Monitoring", title_style))
    story.append(Paragraph("Documentation complete - version presentation professionnelle", h1))
    story.append(Paragraph(f"Date de generation: {DATE_STR}", body))
    story.append(Spacer(1, 14))

    for line in PROJECT_SUMMARY:
        story.append(Paragraph(f"- {line}", body))
    story.append(Spacer(1, 12))

    story.append(Paragraph("1. Architecture de la solution", h1))
    table = Table(ARCH_COMPONENTS, colWidths=[4.0 * cm, 4.2 * cm, 7.2 * cm])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#12355B")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#D0D7E2")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.whitesmoke, colors.HexColor("#F5F9FF")]),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]
        )
    )
    story.append(table)

    story.append(Spacer(1, 12))
    story.append(Paragraph("2. Fonctionnalites metier", h1))
    for item in KEY_FEATURES:
        story.append(Paragraph(f"- {item}", body))

    story.append(PageBreak())

    story.append(Paragraph("3. Cartographie API", h1))
    api_table = Table(API_DOMAINS, colWidths=[5.0 * cm, 10.4 * cm])
    api_table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1E5F74")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#D0D7E2")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.whitesmoke, colors.HexColor("#F8FCFF")]),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]
        )
    )
    story.append(api_table)

    story.append(Spacer(1, 12))
    story.append(Paragraph("4. Deploiement et operations", h1))
    for step in DEPLOY_STEPS:
        story.append(Paragraph(step, body))

    story.append(Spacer(1, 10))
    story.append(Paragraph("5. Agents de collecte", h1))
    for item in AGENT_NOTES:
        story.append(Paragraph(f"- {item}", body))

    story.append(Spacer(1, 10))
    story.append(Paragraph("6. Securite et gouvernance", h1))
    for item in SECURITY_PRACTICES:
        story.append(Paragraph(f"- {item}", body))

    story.append(Spacer(1, 10))
    story.append(Paragraph("7. Roadmap", h1))
    for item in ROADMAP:
        story.append(Paragraph(f"- {item}", body))

    doc.build(story, onFirstPage=header_footer, onLaterPages=header_footer)



def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle



def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)



def add_table_slide(prs, title, data, left=0.6, top=1.6, width=12.2, height=4.8):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(14 if r == 0 else 12)
                if r == 0:
                    para.font.bold = True

    for c in range(cols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(18, 53, 91)
        for para in cell.text_frame.paragraphs:
            para.font.color.rgb = RGBColor(255, 255, 255)



def build_pptx(path: Path) -> None:
    prs = Presentation()

    add_title_slide(
        prs,
        "IT Monitoring",
        "Plateforme de supervision et de gestion de parc IT\nPresentation professionnelle - " + DATE_STR,
    )

    add_bullets_slide(
        prs,
        "Vision et objectifs",
        [
            "Centraliser la visibilite de l'infrastructure IT",
            "Industrialiser la collecte des donnees endpoint",
            "Accelerer la detection et le traitement des incidents",
            "Structurer la gouvernance securite et operations",
        ],
    )

    add_table_slide(prs, "Architecture de reference", ARCH_COMPONENTS)

    add_bullets_slide(
        prs,
        "Capacites fonctionnelles",
        KEY_FEATURES,
    )

    add_table_slide(prs, "Panorama API (extraits)", API_DOMAINS)

    add_bullets_slide(
        prs,
        "Agents et deploiement terrain",
        [
            "Agents disponibles: PowerShell, Python, Go (service Windows)",
            "URL backend par defaut: http://192.168.196.134:8000/assets/scan",
            "Modes de deploiement: manuel, bulk, GPO, all-signed",
            "Surcharge possible via ASSET_BACKEND_URL et config.json",
        ],
    )

    add_bullets_slide(
        prs,
        "Exploitation et fiabilite",
        [
            "Conteneurisation complete: db + backend + frontend",
            "Verification post-deploiement: ports 8000/3000 + health checks HTTP",
            "Logs techniques pour diagnostic rapide en production",
            "Approche progressive vers CI/CD et tests automatises",
        ],
    )

    add_bullets_slide(
        prs,
        "Securite et conformite",
        SECURITY_PRACTICES,
    )

    add_bullets_slide(
        prs,
        "Roadmap professionnelle",
        ROADMAP,
    )

    add_bullets_slide(
        prs,
        "Conclusion",
        [
            "Projet operationnel et presentable a un public IT professionnel",
            "Architecture evolutive orientee exploitation et securite",
            "Base solide pour industrialisation (devops + gouvernance)",
            "Prochaine etape: generalisation du deploiement agent et KPI executifs",
        ],
    )

    prs.save(str(path))



def main() -> None:
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    build_pdf(PDF_PATH)
    build_pptx(PPTX_PATH)
    print(f"PDF generated: {PDF_PATH}")
    print(f"PPTX generated: {PPTX_PATH}")


if __name__ == "__main__":
    main()
